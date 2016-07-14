# -*- coding: utf-8 -*-
"""

Neal Gordon

to be executed directly in bash or cmd.exe
$ python PythonReportTool.py

In order to access the Abaqus odb file, the python script must be run in the abaqus environment, which will use a abaqus license token
abaqus cae noGUI=myscript.py # launches cae and runs script

Inputs: location and name of an OBD file

Outputs: pptx report of file including screenshots, 

"""

import os
import glob
import matplotlib.pyplot as plt
import pandas as pd

# create presentation
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from moviepy.editor import VideoFileClip

def abaqus_cmd(mycmd):
    """ used to execute abaqus commands in the windows OS console
        inputs : mycmd, an abaqus command
        also can use os.system('dir') to run windows or bash commands"""
    import subprocess, sys
    try:
        retcode = subprocess.call(mycmd,shell=True)
        if retcode < 0:
            print >>sys.stderr, mycmd+"...failed during execution", -retcode
        else:
            print >>sys.stderr, mycmd
    except OSError as e:
        print >>sys.stderr, mycmd+"...failed at execution", e

def make_pptx(odbname,dict_o_c,list_v):
    """ Generic Function to create a pptx file
        filename = relative path with directories, eg temp0/beam_tutorial.odb
       creates a new pptx file and inserts all the images in the current directory
       ppt default width 10 inches by 7.5 inches landscape
       #odbname = 'customBeamExample/customBeamExample.odb'
       
        # make a new slide for each component given
        dict_o_c = {'S':['Mises'], 'U':['U1','Magnitude'],  'RF':['RF1'], 'EVOL':['']}
    
        # add all three views onto the same slide
        list_v = ['Iso','Left','Top']      
        
        CBE_make_pptx(odbname,dict_o_c,list_v)
       """

    def img_insert(imgtyp='', picname=''):
        """functions for inserting standard images"""
        try:
            if imgtyp == 'icon':
                left = Inches(0.05) ; top = Inches(6.85)
                slide.shapes.add_picture('https://www.python.org/static/img/python-logo.png', left, top)
            elif imgtyp[:3] == 'Iso':
                left = Inches(6) ; top = Inches(4.25) ; myHeight = Inches(2.5)
                slide.shapes.add_picture(picname, left, top, height=myHeight)
            elif imgtyp == 'Top':
                left = Inches(6) ; top = Inches(2.0) ; myHeight = Inches(2.5)
                slide.shapes.add_picture(picname, left, top, height=myHeight)
            elif imgtyp == 'Left':
                left = Inches(8.2) ; top = Inches(5.8) ; myHeight = Inches(1)
                slide.shapes.add_picture(picname, left, top, height=myHeight)
            elif imgtyp == 'gif':
                left = Inches(0.1) ; top = Inches(2.7) ; myHeight = Inches(4)
                slide.shapes.add_picture(picname, left, top, height=myHeight)
            elif imgtyp == 'main':
                left = Inches(1.9) ; top = Inches(2.5) ; myHeight = Inches(4)
                slide.shapes.add_picture(picname, left, top, height=myHeight)
            else:
                pass
        except:
            print('failed to insert image %s' % picname)

    #odbname = 'customBeamExample/customBeamExample.odb'
    pptx1 = Presentation()
    #blank_slide_layout = pptx1.slide_layouts[6]
    bullet_slide_layout = pptx1.slide_layouts[1]
    title_slide_layout = pptx1.slide_layouts[0]
    title_only_layout = pptx1.slide_layouts[5]

    # first slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "test report generation"
    slide.placeholders[1].text = '''Automatic Report Generation Using Python by '''
    left = Inches(4) ; top = Inches(5.5)
    slide.shapes.add_picture('https://www.python.org/static/img/python-logo.png', left, top)

    # second slide #####################################################################################################
    slide = pptx1.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Simualtion Overview'
    tf = body_shape.text_frame
    tf.text = 'A script was created to run a simulation and generate results'

    p = tf.add_paragraph()
    p.text = 'The script was written in python'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'The script builds the model, sets up the boundary and loading \
    conditions, runs the simulation, extracts data for plotting, and \
    contour plots or model images'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'The same script can be adapted to more \
    complex simulations'
    p.level = 0

    left = top = width = height = Inches(1.5)

    # Third Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    img_insert('icon')
    title_shape.text = 'Underformed Mesh'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Element boundaries have been removed to show the undeformed \n shape of the model"
    picnames = glob.glob(os.path.dirname(odbname)+'/*.png')
    if picnames:
        picname = [f for f in picnames if 'Undeformed' in f]
        if picname:
            img_insert('main', picname[0])

    # Field slides #####################################################################################################
    for o in dict_o_c:
        for c in dict_o_c[o]:
            slide = pptx1.slides.add_slide(title_only_layout)
            title_shape = slide.shapes.title
            img_insert('icon')            
            title_shape.text = o + '-' + c
            txBox = slide.shapes.add_textbox(left, top, width, height)
            txBox.text_frame.text = "This is text inside a textbox\n and this is some more text that you can add"
            for v in list_v:
                picnames = glob.glob(os.path.dirname(odbname)+'/*.png')
                if picnames:
                    picname = [f for f in picnames if o in f if c in f if v in f]
                    if picname:
                        img_insert(v, picname[0])
            picnames = glob.glob(os.path.dirname(odbname)+'/*.gif')
            if picnames:
                picname = [f for f in picnames if o in f if c in f if 'Iso' in f]
                if picname:
                    img_insert('gif', picname[0])

    # save pptx
    pptx1.save(os.path.splitext(odbname)[0]+'.pptx')

def HS_make_pptx(odbname):
    """ Fucntion specific to Heat-Set report"""

    def insert_img(f='', o='',c='', v='', imtype='png', left=0, top=0, height=1):
        # o='STH',c='', v='Back',
        picnames = glob.glob(os.path.dirname(odbname)+'/*.'+imtype)
        if picnames:
            picname = [fn for fn in picnames if f in fn if o in fn if c in fn if v in fn]
            if picname:
                if height:
                    slide.shapes.add_picture(picname[0], Inches(left), Inches(top), height=Inches(height))
                else: # choose default
                    slide.shapes.add_picture(picname[0], Inches(left), Inches(top))

    def insert_icon():
        slide.shapes.add_picture('https://www.python.org/static/img/python-logo.png', Inches(0.05) , Inches(6.85))

    #odbname = 'customBeamExample/customBeamExample.odb'
    pptx1 = Presentation()
    #blank_slide_layout = pptx1.slide_layouts[6]
    bullet_slide_layout = pptx1.slide_layouts[1]
    title_slide_layout = pptx1.slide_layouts[0]
    title_only_layout = pptx1.slide_layouts[5]

    # first slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Project Name"
    slide.placeholders[1].text = 'Project Number: XXXXXX\nDrawing Number: XXXXXXXX REV X\nWeight: XX grams\nPreform Number: XXXXXXXX REV X'
    insert_icon()
    insert_img('Undeformed','','','','png',-.25, 0, 3)

    # second slide #####################################################################################################
    slide = pptx1.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Simulation Characteristics'
    tf = body_shape.text_frame

    tf.text = 'Heat Set'
    p = tf.add_paragraph()
    p.text = '''Fill Line: 1.5 in'''
    p.level = 1
    p = tf.add_paragraph()
    p.text = '''Fill Temp: 185 °F'''
    p.level = 1
    p = tf.add_paragraph()
    p.text = '''Room Temp: 75 °F'''
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Vacuum Resistance'
    p.level = 0
    p = tf.add_paragraph()
    p.text = '''Volume of Water Removed: 5.0 mL'''
    p.level = 1
    p = tf.add_paragraph()
    p.text = '''Volume of Headspace Removed: 1.0 mL'''
    p.level = 1

    # Third Slide #####################################################################################################
    left = top = width = height = Inches(1.5)
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Thickness Profile'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Determinate of Profile: estimated with preform xxxxxxx"
    insert_icon()
    insert_img('frame-1','STH','','Back',  'png', 5.5,    2, 5)
    insert_img('frame-1','STH','','Bottom','png', 5.25, 2, 5)

    # Fourth Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Deformation After Heat Set'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Max Expected Deformation: 0.078 in\nLocation of Max Deformation: vacuum panels"
    insert_icon()
    insert_img('',       'U','U1','Iso1',  'gif', 0.1,  2.7, 4)
    insert_img('frame-1','U','U1','Back',  'png', 6,    5.3, 2)
    insert_img('frame-1','U','U1','Bottom','png', 6.25, 2.5, 3)

    # Fifth Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Stress After Heat Set'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Max Expected Stress: 6126 psi"
    insert_icon()
    insert_img('',       'S','S11','Iso1',  'gif', 0.1,  2.7, 4)
    insert_img('frame-1','S','S11','Back',  'png', 6,    5.3, 2)
    insert_img('frame-1','S','S11','Bottom','png', 6.25, 2.5, 3)

    # Sixth Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Deformation After Headspace Removed'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Max Expected Deformation: 0.081 in\nLocation of Max Deformation: vacuum panels"
    insert_icon()

    # Seventh Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Deformation After Volume of Water Removed'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Max Expected Deformation: .088 in\nLocation of Max Deformation: vacuum panels"
    insert_img(o='U',c='U11')
    insert_icon()

    # save pptx
    pptx1.save(os.path.splitext(odbname)[0]+'.pptx')

def BT_make_pptx(odbname):
    """ Fucntion specific to Heat-Set report"""

    def insert_img(f='frame-1', o='',c='', v='', imtype='png', left=0, top=0, height=0):
        # o='STH',c='', v='Back',
        picnames = glob.glob(os.path.dirname(odbname)+'/*.'+imtype)
        if picnames:
            picname = [fn for fn in picnames if f in fn if o in fn if c in fn if v in fn]
            if picname:
                if height:
                    slide.shapes.add_picture(picname[0], Inches(left), Inches(top), height=Inches(height))
                else: # choose default
                    slide.shapes.add_picture(picname[0], Inches(left), Inches(top))

    def insert_icon():
        slide.shapes.add_picture('https://www.python.org/static/img/python-logo.png', Inches(0.05) , Inches(6.85))

    #odbname = 'customBeamExample/customBeamExample.odb'
    pptx1 = Presentation()
    #blank_slide_layout = pptx1.slide_layouts[6]
    bullet_slide_layout = pptx1.slide_layouts[1]
    title_slide_layout = pptx1.slide_layouts[0]
    title_only_layout = pptx1.slide_layouts[5]

    # first slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Project Name"
    slide.placeholders[1].text = 'Project Number: XXXXXX\nDrawing Number: XXXXXXXX REV X\nWeight: XX grams\nPreform Number: XXXXXXXX REV X'
    insert_icon()
    insert_img('Undeformed','','','','png',-.25, 0, 3)

    # second slide #####################################################################################################
    slide = pptx1.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Simulation Characteristics'
    tf = body_shape.text_frame

    tf.text = 'Heat Set'
    p = tf.add_paragraph()
    p.text = '''Fill Line: xx in'''
    p.level = 1
    p = tf.add_paragraph()
    p.text = '''Fill Temp: xx °F'''
    p.level = 1
    p = tf.add_paragraph()
    p.text = '''Room Temp: xx °F'''
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Vacuum Resistance'
    p.level = 0
    p = tf.add_paragraph()
    p.text = '''Volume of Water Removed: xx mL'''
    p.level = 1
    p = tf.add_paragraph()
    p.text = '''Volume of Headspace Removed: xx mL'''
    p.level = 1

    # Third Slide #####################################################################################################
    left = top = width = height = Inches(1.5)
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Thickness Profile'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Determinate of Profile: estimated with preform xxxxxxx"
    insert_icon()
    insert_img('','STH','','','png', 5.25, 2, 5)

    # Fourth Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Deformation After Heat Set'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Max Expected Deformation: xxx in\nLocation of Max Deformation: vacuum panels"
    insert_icon()
    insert_img('','U','Magnitude','','png', 5.25, 2, 5)
    insert_img('','U','Magnitude','','gif', 0.25, 2.25, 4)

    # Fifth Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Stress After Heat Set'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Max Expected Stress: xxx psi"
    insert_icon()
    insert_img('','S','Mises','','png', 5.25, 2, 5)
    insert_img('','S','Mises','','gif', 0.25, 2.25, 4)

    # Sixth Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Deformation After Headspace Removed'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Max Expected Deformation:xxx in\nLocation of Max Deformation: vacuum panels"
    insert_icon()
    insert_img('','U','CSYS-CYN','','png', 5.25, 2, 3.5)
    insert_img('','U','CSYS-CYN','','gif', 0.25, 2.25, 3.5)

    # Seventh Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Plots of data'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "Reaction Force and Displacement"
    insert_img('RF1','U1','XYplot','','png', 0.5, 2, 4)
    insert_img('Xplot','RF1','','','png', 5, 2, 4)
    insert_icon()

    # Eighth Slide #####################################################################################################
    slide = pptx1.slides.add_slide(title_only_layout)
    title_shape = slide.shapes.title
    title_shape.text = 'Mises Stress with plot'
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text = "plot Reaction force with with Displacement"
    insert_img('RF1','U1','XYplot','','gif', 5, 2, 3.5)
    insert_img('S','Mises','Iso1','CSYS-CYN','gif', 0.5, 2, 3.5)
    insert_icon()

    # save pptx
    pptx1.save(os.path.splitext(odbname)[0]+'.pptx')

def avi_to_gif(odbname):
    """converts an avi video file to an animated gif so that it can easily be
        easily inserted into a pptx
        avidir = 'customBeamExample'
    """
    avinames=glob.glob(os.path.dirname(odbname)+'/*.avi')
    for aviname in avinames:
        clip = VideoFileClip(aviname)
        clip.write_gif(os.path.splitext(aviname)[0]+'.gif')
        #os.remove(aviname)

def field_plot(y):
    '''function designed to import a csv with a single varialbe and plot it and save a png '''

    plt.rcParams['font.size'] = 14

    df_y = pd.read_csv(y,names=[y],delimiter=',')
    #print df
    print y
    plt.figure(1, figsize=(10,6))
    plt.plot(df_y)
    plt.xlabel('strain')
    plt.ylabel('stress')
    plt.legend('')
    plt.title(y)
    plt.savefig(os.path.splitext(y)[0]+'plot.png')

def history_plots(csvnames):
    """
    'customBeamExample/custom_beam_tutorial_beamLoad_Element BEAMINSTANCE.42 Int Point 1_E11.csv'
     after running odb_to_csv, this function can read those files and plot

     a simple way to check the contents of csvs print(open(filename1).read())
     """
     
    for csvname in csvnames:
        df = pd.read_csv(csvname,names=['time','val'],delimiter=',')
        #print df
        plt.figure()
        plt.plot(df.time, df.val)
        plt.xlabel('time')
        #plt.ylabel(csvname)
        plt.title(csvname)

def read_rpt(rptname):
    """ Reads a rpt file generated by abaqus and can plot the contents
        rptname = "C:/Users/ngordon/Desktop/abaqus-scripting/customBeamExample/customBeamExample_step-0_frame-0_field-E_elsets-_csys-.rpt"
    """
    with open(rptname) as f:
       for r, line in enumerate(f):
           if 'Element' in line or 'Nodal' in line:
               if line.split()[0] == 'Element' or line.split()[0] == 'Nodal':
                   columnLabels=line.split()
                   break
    df = pd.read_csv(rptname,skiprows= r+3, delim_whitespace=True, index_col=False, names=columnLabels)
    #df.head(2)

    return df

def get_rpts(odbname):
    '''show the contens of all the rpt files in a directory'''
    rptnames = glob.glob(os.path.dirname(odbname)+'/*.rpt')
    for rptname in rptnames:
        df = read_rpt(rptname)
        #df['S.Mises'].plot()
        print df.head(2)

def image_resize(picnames=[]):
    """ This program does not use this but may be needed
        picnames ia a list of the full path of images
        do not recommend using this, change the picture size in the AbaqusReportTool function by adjusting
            AbaqusReportTool.VP_edit(myView='Iso', vp_width=225, replaceKeyword = '')
    """
    for picname in picnames:
        im1 = Image.open(picname)
        imgscale = 0.5
        width,height = im1.size
        newsize = (int(width*imgscale), int(height*imgscale))
        im1 = im1.resize(newsize, Image.ANTIALIAS)
        os.remove(picname)
        im1.save(picname)

def clearFiles(odbname, filetypes = ['png','csv','txt','avi','gif']):
    """clears the files """
    for fileext in filetypes:
        for f in glob.glob(os.path.dirname(odbname)+'/*.'+fileext):
            try:
                os.remove(f)
            except:
                pass

def clearFilesExcept(odbname, filetypes = ['.odb','.py','.txt','.pptx','.cae','.stp','CATPart']):
    """clears the files except filetypes"""
    for f in glob.glob(os.path.dirname(odbname)+'/*'):
        if os.path.splitext(f)[1] not in filetypes:
            try:
                os.remove(f)
            except:
                pass

def makeReportData(reportType, odbname):
    """generic report generation with a flag indicating what type of report it is
    abaqus cae noGUI=AbaqusReportTool.py -- BT bottle_test/bottle_test.odb
    """
    abaqus_cmd('abaqus cae noGUI=AbaqusReportTool.py -- '+reportType+' '+odbname)

def make_dict_from_files(odbname):
    """ this will generate a dictioanry from the png files in directory to use for report generation."""
    picnames = glob.glob(os.path.dirname(odbname)+'/*.png')
    list_v = []
    dict_o_c = {}
    for f in picnames:
        s = os.path.basename(os.path.splitext(f)[0]).split('_')
        if s[-1] not in list_v:
            list_v.append(s[-1])
        if s[-2] != 'Undeformed':
            if s[-3] in dict_o_c.keys():
                if s[-2] not in dict_o_c[s[-3]]:
                    dict_o_c[s[-3]].append(s[-2])
            else:
                dict_o_c[s[-3]] = [s[-2]]
    return list_v, dict_o_c

# Main functions #####################################################################################################
def beamExample():
    """ demo of  executing abaqus commands """
    abaqus_cmd('abaqus cae noGUI=AbaqusReportTool.py -- beamExample')
    odbname = 'beamExample/beam_tutorial.odb'
    makeReportData('BE', odbname)
    avi_to_gif(odbname)
    list_v, dict_o_c = make_dict_from_files(odbname)
    make_pptx(odbname,dict_o_c,list_v)

def customBeamExample():
    """ demo created of a custom abaqus example that generates a stress-strain plot"""
    abaqus_cmd('abaqus cae noGUI=AbaqusReportTool.py -- customBeamExample')
    odbname = 'customBeamExample/customBeamExample.odb'
    makeReportData('CBE', odbname)
    avi_to_gif(odbname)
    list_v, dict_o_c = make_dict_from_files(odbname)
    make_pptx(odbname,dict_o_c,list_v)

def bottle_test():
    odbname = 'bottle_test/bottle_test.odb'
    if not glob.glob('bottle_test/*.odb'):
        abaqus_cmd('abaqus cae noGUI=AbaqusReportTool.py -- bottle_test')
    makeReportData('BT', odbname)
    avi_to_gif(odbname)
    BT_make_pptx(odbname)
    clearFilesExcept(odbname)

def Reports():
    """ will look in the folders for odbfiles and attempt to generate a report for every odb file there"""
    # find the name of the odbfile in the in HS_Example
    for d in ['HS_Example','TL_Example']:
        odbnames = glob.glob('Reports/'+d+'/*.odb')
        if odbnames:
            for odbname in odbnames:
                try:
                    if ' ' in odbname or set('[~!@$%^*+{}":;]').intersection(odbname):
                        print('%s \n  has spaces or special characters in the filename and or filepath, error...\n' % odbname)
                        continue
                    else:
                        print('%s \n  running makeReportData, OK...\n' % odbname)
                        makeReportData(d[:2], odbname)
                        avi_to_gif(odbname)
                        HS_make_pptx(odbname)
                        clearFilesExcept(odbname)
                        print(' ')
                except:
                    print('failed creating report for %s' % odbname)

if __name__ == '__main__':
    """ if this py file is executed these statements will be executed"""
    #beamExample()
    #customBeamExample()
    bottle_test()
    #Reports()
